import matplotlib
matplotlib.use("Agg")
import json
import pandas as pd
from Api import BaseAPI
import statistics
from numerize import numerize
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE,PP_ALIGN,MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR,MSO_SHAPE
import matplotlib.colors as Colour
from utils.create_table import Table
from  utils.create_heading import Heading
from datetime import datetime,timedelta
import matplotlib.pyplot as plt
from matplotlib.dates import DateFormatter
import uuid
import os
class BusinessDescription(BaseAPI):
    __url__1="https://yh-finance.p.rapidapi.com/stock/v2/get-profile"
    __url__2="https://yh-finance.p.rapidapi.com/stock/v2/get-financials"
    __url__3="https://yh-finance.p.rapidapi.com/stock/v2/get-cash-flow"
    __url__4="https://yh-finance.p.rapidapi.com/stock/v3/get-historical-data"
    @classmethod
    def get_business_descripton(cls,raw_data):
        data:str=cls.search(raw_data,"longBusinessSummary")
        old_sentences=data.split(".")
        new_sentences=[]
        for idx in range(len(old_sentences)):
            if not old_sentences[idx].endswith("Inc"):
                new_sentence=old_sentences[idx].rstrip()+".\n"
            else:
                new_sentence=old_sentences[idx].rstrip()+"."
            if idx!=len(old_sentences)-1:
                new_sentences.append(new_sentence)
            else:
                new_sentences.append(old_sentences[idx]+".")
        text="".join(new_sentences)
        return text
    @classmethod
    def get_financial_summary(cls,raw_data_financial,raw_data_cashflow):
        column_data={}
        revenue=cls.search(raw_data_financial,"annualTotalRevenue")
        ebitda=cls.search(raw_data_financial,"annualEbitda")
        fcf=cls.search(raw_data_cashflow,"annualFreeCashFlow")
        prev_rev=revenue[-4]["reportedValue"]["raw"]
        for dataPoint in range(-3,0):
            year=datetime.strptime(revenue[dataPoint]["asOfDate"],"%Y-%m-%d").year
            key=str(year)
            cur_rev=revenue[dataPoint]["reportedValue"]["raw"]
            if key not in column_data:
                column_data[key]=[]
            growth=(cur_rev-prev_rev)/prev_rev
            cur_ebitda=ebitda[dataPoint]["reportedValue"]["raw"]
            ebitda_margin=cur_ebitda/cur_rev
            cur_fcf=fcf[dataPoint]["reportedValue"]["raw"]
            fcf_margin=cur_fcf/cur_rev
            column_data[key].append(numerize.numerize(cur_rev))
            column_data[key].append(numerize.numerize(growth*100)+"%")
            column_data[key].append(numerize.numerize(cur_ebitda))
            column_data[key].append(numerize.numerize(ebitda_margin)+"%")
            column_data[key].append(numerize.numerize(cur_fcf))
            column_data[key].append(numerize.numerize(fcf_margin)+"%")
            prev_rev=cur_rev
        indexes=["Revenue","Growth(%)","EBITDA","Margin(%)","FCF","Margin(%)"]
        return {"indexes":indexes,"columns":column_data}
    @classmethod
    def get_share_price_performance(cls,raw_data):
        df=pd.json_normalize(raw_data["prices"])
        Y=df["close"]
        df["date"]=pd.to_datetime(df["date"],unit="s")
        first_date=df["date"][0]
        ticks=[]
        ticks.append(first_date)
        for i in df["date"]:
            first_date=first_date-timedelta(days=30)
            ticks.append(first_date)
        figure=plt.figure(figsize=(6,3))
        fig=plt.plot(df["date"],Y,label='Close',color='#ffa500')
        plt.xticks=ticks
        plt.gca().xaxis.set_major_formatter(DateFormatter("%b,%Y"))
        plt.legend()
        plt.gcf().autofmt_xdate()
        for tick in plt.gca().xaxis.get_major_ticks():
            tick.label.set_fontsize(9)
        for tick in plt.gca().yaxis.get_major_ticks():
            tick.label.set_fontsize(9)
        plt.grid(True)
        plt.gcf().subplots_adjust(top=1,bottom=0.25)
        filename=uuid.uuid4().hex[:8].upper()
        path=os.path.join(os.getcwd(),'Graphs')
        if  not os.path.exists(path):
            os.mkdir(path)
        adress=path+'\\'+filename+'.png'
        plt.savefig(adress)
        return adress
    @classmethod
    def ppt_output(cls,pr,company,style):
        try:
            HEADING_FONT=style["heading-font"]
        except:
            HEADING_FONT="Calibri"
        try:
            HEADING_FONT_SIZE=int(style["heading-font-size"])
        except:
            HEADING_FONT_SIZE=35
        try:
            SUB_HEADING_FONT=style["sub-heading-font"]
        except:
            SUB_HEADING_FONT="Calibri"
        try:
            BODY_TEXT_FONT=style["body-text-font"]
        except:
            BODY_TEXT_FONT="Calibri"
        try:
            BODY_TEXT_FONT_SIZE=style["body-text-font-size"]
        except:
            BODY_TEXT_FONT_SIZE=13
        BODY_TEXT_HEADING_FONT_SIZE=BODY_TEXT_FONT_SIZE+1
        try:
            THEME_COLOR=tuple(int(255*i) for i in Colour.to_rgb(style["theme-color"]))
        except:
            THEME_COLOR=(1,39,99)
        #  Adding slide
        layout=pr.slide_layouts[6] # Blank layout
        slide=pr.slides.add_slide(layout)
        # # Heading
        heading=Heading(HEADING_FONT_SIZE,THEME_COLOR,HEADING_FONT,(Inches(0.2),Inches(0),Inches(8),Inches(0.5)),"Company overview",slide)
        heading.create_heading()
        # Separation line
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0),Inches(0.7),Inches(2.1),Inches(0.7))
        fill=line.line
        fill.color.rgb=RGBColor(184,25,4)
        # Heading
        bd_heading_shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.5),Inches(1),Inches(10),Inches(0.3))
        bd_heading_txtFrame=bd_heading_shape.text_frame
        bd_heading_txtFrame.text="Buisness Description"
        try:
            bd_heading_txtFrame.fit_text(font_family=str(SUB_HEADING_FONT))
        except:
            bd_heading_txtFrame.fit_text(font_family=u"Calibri")
        bd_shape_fill=bd_heading_shape.fill
        bd_shape_fill.solid()
        bd_shape_fill.fore_color.rgb=RGBColor(*THEME_COLOR)
        #content
        bd_data=cls.base_fetch({"symbol":company,"region":"US"},url=cls.__url__1)
        bd_content_txtBox=slide.shapes.add_textbox(Inches(0.5),Inches(1.5),Inches(10),Inches(6))
        bd_content_txtFrame=bd_content_txtBox.text_frame
        bd_content_txtFrame.word_wrap=True
        bd_content_txtFrame.text=cls.get_business_descripton(bd_data)
        # Heading
        bd_heading_shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(11),Inches(1),Inches(6.5),Inches(0.3))
        bd_heading_txtFrame=bd_heading_shape.text_frame
        bd_heading_txtFrame.text="Financial summary"
        try:
            bd_heading_txtFrame.fit_text(font_family=str(SUB_HEADING_FONT))
        except:
            bd_heading_txtFrame.fit_text(font_family=u"Calibri")
        bd_shape_fill=bd_heading_shape.fill
        bd_shape_fill.solid()
        bd_shape_fill.fore_color.rgb=RGBColor(*THEME_COLOR)
        #table
        raw_data_financials=cls.base_fetch({"symbol":company,"region":"US"},url=cls.__url__2)
        raw_data_fcf=cls.base_fetch({"symbol":company,"region":"US"},url=cls.__url__3)
        fs_data=cls.get_financial_summary(raw_data_financials,raw_data_fcf)
        column_width={"normal":1
                      ,"index":1.5}
        table=Table(fs_data,(12,1.5,len(fs_data["columns"])*1+1.5,len(fs_data["indexes"])*0.2),column_width,slide)
        table.create()
        # SPP
        # Heading
        bd_heading_shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(11),Inches(4),Inches(6.5),Inches(0.3))
        bd_heading_txtFrame=bd_heading_shape.text_frame
        bd_heading_txtFrame.text="Share Price Performance"
        try:
            bd_heading_txtFrame.fit_text(font_family=str(SUB_HEADING_FONT))
        except:
            bd_heading_txtFrame.fit_text(font_family=u"Calibri")
        bd_shape_fill=bd_heading_shape.fill
        bd_shape_fill.solid()
        bd_shape_fill.fore_color.rgb=RGBColor(*THEME_COLOR)
        sp_data=cls.base_fetch({"symbol":company,"region":"US"},cls.__url__4)
        adress=cls.get_share_price_performance(sp_data)
        graph=slide.shapes.add_picture(adress,Inches(11),Inches(4.5),Inches(6.5),Inches(3))
        os.remove(adress)
        return
        
        
        
        
        

            
        
    