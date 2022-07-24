from pptx.util import Inches,Pt
import json
from pptx.enum.text import PP_ALIGN
class Table:
    def __init__(self,data:dict,position:tuple,column_width:dict,slide):
        self.data=data
        self.position=position
        self.slide=slide
        self.column_width=column_width
    def create(self):
        table=self.slide.shapes.add_table(len(self.data["indexes"])+1,len(self.data["columns"])+1,Inches(self.position[0]),Inches(self.position[1]),Inches(self.position[2]),Inches(self.position[3])).table
        for column in range(len(self.data["columns"])+1):
            paragraph=table.cell(0,column).text_frame.paragraphs[0]
            if column==0:
                paragraph.text="USD"
            else:
                # paragraph.text=list(self.data["columns"].keys())[len(self.data["columns"])-column]
                paragraph.text=list(self.data["columns"].keys())[column-1]
            paragraph.font.size=Pt(12)
            paragraph.alignment=PP_ALIGN.LEFT
            table.columns[column].width=Inches(self.column_width["normal"])
        table.columns[0].width=Inches(self.column_width["index"])
        for index in range(len(self.data["indexes"])):
            paragraph=table.cell(index+1,0).text_frame.paragraphs[0]
            paragraph.text=self.data["indexes"][index]
            paragraph.font.size=Pt(12)
            paragraph.alignment=PP_ALIGN.LEFT
        column_index=1
        for value in self.data["columns"].values():
            for data_index in range(len(value)):
                paragraph=table.cell(data_index+1,column_index).text_frame.paragraphs[0]
                paragraph.text=value[data_index] if value[data_index] is not None else "NA"
                paragraph.font.size=Pt(12)
                paragraph.alignment=PP_ALIGN.CENTER 
            column_index+=1