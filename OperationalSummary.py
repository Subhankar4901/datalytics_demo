import json
import sys
from Api import BaseAPI
import statistics
from numerize import numerize
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE,PP_ALIGN,MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR,MSO_SHAPE
import matplotlib.colors as Colour
from utils.create_table import Table
from  utils.create_heading import Heading
from datetime import datetime
class OperationalSummary(BaseAPI):
    __url__1="https://yh-finance.p.rapidapi.com/stock/v2/get-statistics"
    __url__2="https://yh-finance.p.rapidapi.com/stock/v2/get-financials"
    __url__3="https://yh-finance.p.rapidapi.com/stock/v2/get-cash-flow"
    @classmethod
    def set_stock_price(cls,raw_data,company_data):
        if "Stock Price" not in company_data:
            company_data["Stock Price"]=[]
        stock_price=cls.search(raw_data,"regularMarketPrice")
        if stock_price is not None:
            company_data["Stock Price"].append(stock_price["fmt"])
        else:
            company_data["Stock Price"].append(None)
    @classmethod
    def set_enterprise_value(cls,raw_data,company_data):
        if "Enterprise Value" not in company_data:
            company_data["Enterprise Value"]=[]
        enter_price=cls.search(raw_data,"enterpriseValue")
        if enter_price is not None:
            company_data["Enterprise Value"].append(enter_price["fmt"])
        else:
            company_data["Enterprise Value"].append(None)
    @classmethod
    def set_revenue(cls,raw_data,company_data):
        if "Rev. CAGR" not in company_data:
            company_data["Rev. CAGR"]=[]
        rev=cls.search(raw_data,"annualTotalRevenue")
        first=rev[-3]
        last=rev[-1]
        try:
            cagr=((last["reportedValue"]["raw"]/first["reportedValue"]["raw"])**(1/3))-1
            cagr=numerize.numerize(cagr)
        except:
            cagr=None
        company_data["Rev. CAGR"].append(cagr)
    @classmethod
    def set_EBITDA_margin(cls,raw_data,company_data):
        ebitda=cls.search(raw_data,"annualEbitda")
        revenue=cls.search(raw_data,"annualTotalRevenue")
        for datapoint in range(-3,0):
            try:
                year=datetime.strptime(ebitda[datapoint]["asOfDate"],"%Y-%m-%d").year
                key="EBITDA Margin\n"+str(year)
                if key not in company_data:
                    company_data[key]=[]
                ebitda_margin=ebitda[datapoint]["reportedValue"]["raw"]/revenue[datapoint]["reportedValue"]["raw"]
                ebitda_margin=str(numerize.numerize(ebitda_margin))+"%"
                company_data[key].append(ebitda_margin)
            except:
                if datapoint==-1:
                    year=str(datetime.now().year-1)
                if datapoint==-2:
                    year=str(datetime.now().year-2)
                if datapoint==-3:
                    year=str(datetime.now().year-3)
                key="EBITDA Margin\n"+str(year)
                if key not in company_data:
                    company_data[key]=[]
                company_data[key].append(None)
    @classmethod
    def set_EBIT_margin(cls,raw_data,company_data):
        ebit=cls.search(raw_data,"annualOperatingIncome")
        revenue=cls.search(raw_data,"annualTotalRevenue")
        for datapoint in range(-3,0):
            try:
                year=datetime.strptime(ebit[datapoint]["asOfDate"],"%Y-%m-%d").year
                key="EBIT Margin\n"+str(year)
                if key not in company_data:
                    company_data[key]=[]
                ebit_margin=ebit[datapoint]["reportedValue"]["raw"]/revenue[datapoint]["reportedValue"]["raw"]
                ebit_margin=str(numerize.numerize(ebit_margin))+"%"
                company_data[key].append(ebit_margin)
            except:
                if datapoint==-1:
                    year=str(datetime.now().year-1)
                if datapoint==-2:
                    year=str(datetime.now().year-2)
                if datapoint==-3:
                    year=str(datetime.now().year-3)
                key="EBIT Margin\n"+str(year)
                if key not in company_data:
                    company_data[key]=[]
                company_data[key].append(None)
                        
    @classmethod
    def set_FCF_margin(cls,raw_data_fcf,raw_data_statistics,company_data):
        fcf=cls.search(raw_data_fcf,"annualFreeCashFlow")
        revenue=cls.search(raw_data_statistics,"annualTotalRevenue")
        for datapoint in range(-3,0):
            try:
                year=datetime.strptime(fcf[datapoint]["asOfDate"],"%Y-%m-%d").year
                key="FCF Margin\n"+str(year)
                if key not in company_data:
                    company_data[key]=[]
                fcf_margin=fcf[datapoint]["reportedValue"]["raw"]/revenue[datapoint]["reportedValue"]["raw"]
                fcf_margin=str(numerize.numerize(fcf_margin))+"%"
                company_data[key].append(fcf_margin)
            except:
                if datapoint==-1:
                    year=str(datetime.now().year-1)
                if datapoint==-2:
                    year=str(datetime.now().year-2)
                if datapoint==-3:
                    year=str(datetime.now().year-3)
                key="FCF Margin\n"+str(year)
                if key not in company_data:
                    company_data[key]=[]
                company_data[key].append(None)
            
    @classmethod
    def set_netIncome_margin(cls,raw_data,company_data):
        income=cls.search(raw_data,"annualNetIncome")
        revenue=cls.search(raw_data,"annualTotalRevenue")
        for datapoint in range(-3,0):
            try:
                year=datetime.strptime(income[datapoint]["asOfDate"],"%Y-%m-%d").year
                key="Net Income Margin\n"+str(year)
                if key not in company_data:
                    company_data[key]=[]
                income_margin=income[datapoint]["reportedValue"]["raw"]/revenue[datapoint]["reportedValue"]["raw"]
                income_margin=str(numerize.numerize(income_margin))+"%"
                company_data[key].append(income_margin)
            except:
                if datapoint==-1:
                    year=str(datetime.now().year-1)
                if datapoint==-2:
                    year=str(datetime.now().year-2)
                if datapoint==-3:
                    year=str(datetime.now().year-3)
                key="Net Income Margin\n"+str(year)
                if key not in company_data:
                    company_data[key]=[]
                company_data[key].append(None)        
    @classmethod
    def output(cls,companies:list):
        column_data={}
        for company in companies:
            raw_data_statistics=cls.base_fetch({"symbol":company,"region":"US"},cls.__url__1)
            raw_data_financials=cls.base_fetch({"symbol":company,"region":"US"},cls.__url__2)
            raw_data_cash_flow=cls.base_fetch({"symbol":company,"region":"US"},cls.__url__3)
            cls.set_stock_price(raw_data_statistics,column_data)
            cls.set_enterprise_value(raw_data_statistics,column_data)
            cls.set_revenue(raw_data_financials,column_data)
            cls.set_EBITDA_margin(raw_data_financials,column_data)
            cls.set_EBIT_margin(raw_data_financials,column_data)
            cls.set_FCF_margin(raw_data_cash_flow,raw_data_financials,column_data)
            cls.set_netIncome_margin(raw_data_financials,column_data)
        indexes=companies
        return {"columns":column_data,"indexes": indexes}
    @classmethod
    def ppt_output(cls,pr,companies:list,style):
        data=cls.output(companies)
        try:
            HEADING_FONT=style["heading-font"]
        except:
            HEADING_FONT="Calibri"
        try:
            HEADING_FONT_SIZE=int(style["heading-font-size"])
        except:
            HEADING_FONT_SIZE=35
        try:
            SUB_HEADING_FONT=style["sub-heading-font"]
        except:
            SUB_HEADING_FONT="Calibri"
        try:
            BODY_TEXT_FONT=style["body-text-font"]
        except:
            BODY_TEXT_FONT="Calibri"
        try:
            BODY_TEXT_FONT_SIZE=style["body-text-font-size"]
        except:
            BODY_TEXT_FONT_SIZE=13
        BODY_TEXT_HEADING_FONT_SIZE=BODY_TEXT_FONT_SIZE+1
        try:
            THEME_COLOR=tuple(int(255*i) for i in Colour.to_rgb(style["theme-color"]))
        except:
            THEME_COLOR=(1,39,99)
        #  Adding slide
        layout=pr.slide_layouts[6] # Blank layout
        slide=pr.slides.add_slide(layout)
        # # Heading
        heading=Heading(HEADING_FONT_SIZE,THEME_COLOR,HEADING_FONT,(Inches(0.2),Inches(0),Inches(8),Inches(0.5)),"Operational Summary",slide)
        heading.create_heading()
        # Separation line
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0),Inches(0.7),Inches(2.1),Inches(0.7))
        fill=line.line
        fill.color.rgb=RGBColor(184,25,4)
        # Table
        column_width={"normal":0.9
                      ,"index":0.9}
        table=Table(data,(1.5,2.2,len(data["columns"])*0.9+0.9,len(data["indexes"])*0.5),column_width,slide)
        table.create()
        return