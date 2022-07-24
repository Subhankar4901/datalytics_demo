import json
import sys
from Api import BaseAPI
import statistics
from numerize import numerize
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE,PP_ALIGN,MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR,MSO_SHAPE
import matplotlib.colors as Colour
from utils.create_table import Table
from  utils.create_heading import Heading
from datetime import datetime
import matplotlib.pyplot as plt
import os
import uuid
class ValuationOutputs(BaseAPI):
    __url__1="https://yh-finance.p.rapidapi.com/stock/v2/get-financials"
    __url__2="https://yh-finance.p.rapidapi.com/stock/v2/get-statistics"
    __url__3="https://yh-finance.p.rapidapi.com/stock/v2/get-cash-flow"
    @classmethod
    def EVbySales(cls,raw_data_stat,raw_data_fin):
        ev=cls.search(raw_data_stat,"enterpriseValue")
        rev=cls.search(raw_data_fin,"annualTotalRevenue")
        evbysales=[]
        years=[]
        for dataPoint in rev:
            date=str(datetime.strptime(dataPoint["asOfDate"],"%Y-%m-%d").year)
            years.append(date)
            evbysales.append(int(ev["raw"]/dataPoint["reportedValue"]["raw"]))
        fig=plt.figure(figsize=(6,4))
        bar=plt.bar(years,evbysales,0.3,color='#132E57')
        yticks=[int(i) for i in evbysales]
        #plt.gca().set_yticks(yticks)
        plt.gca().spines["left"].set_visible(False)
        plt.gca().spines["right"].set_visible(False)
        plt.gca().spines["top"].set_visible(False)
        plt.gca().spines["bottom"].set_color("#DDDDDD")
        for rect in bar:
            height=rect.get_height()
            plt.text(rect.get_x()+0.15,height/2,str(height)+"x",ha='center', va='bottom')
        #plt.gcf().subplots_adjust(top=1,bottom=0.1,left=0.1)
        plt.gcf().subplots_adjust(top=1,bottom=0.25)
        filename=uuid.uuid4().hex[:8].upper()
        path=os.path.join(os.getcwd(),'Graphs')
        if  not os.path.exists(path):
            os.mkdir(path)
        adress=path+'\\'+filename+'.png'
        plt.savefig(adress)
        return adress
    @classmethod
    def EVbyEbit(cls,raw_data_stat,raw_data_fin):
        ev=cls.search(raw_data_stat,"enterpriseValue")
        rev=cls.search(raw_data_fin,"annualOperatingIncome")
        evbysales=[]
        years=[]
        for dataPoint in rev:
            date=str(datetime.strptime(dataPoint["asOfDate"],"%Y-%m-%d").year)
            years.append(date)
            evbysales.append(int(ev["raw"]/dataPoint["reportedValue"]["raw"]))
        fig=plt.figure(figsize=(6,4))
        bar=plt.bar(years,evbysales,0.3,color='#132E57')
        yticks=[int(i) for i in evbysales]
        #plt.gca().set_yticks(yticks)
        plt.gca().spines["left"].set_visible(False)
        plt.gca().spines["right"].set_visible(False)
        plt.gca().spines["top"].set_visible(False)
        plt.gca().spines["bottom"].set_color("#DDDDDD")
        for rect in bar:
            height=rect.get_height()
            plt.text(rect.get_x()+0.15,height/2,str(height)+"x",ha='center', va='bottom')
        #plt.gcf().subplots_adjust(top=1,bottom=0.1,left=0.1)
        plt.gcf().subplots_adjust(top=1,bottom=0.25)
        filename=uuid.uuid4().hex[:8].upper()
        path=os.path.join(os.getcwd(),'Graphs')
        if  not os.path.exists(path):
            os.mkdir(path)
        adress=path+'\\'+filename+'.png'
        plt.savefig(adress)
        return adress
    @classmethod
    def EVbyEbitda(cls,raw_data_stat,raw_data_fin):
        ev=cls.search(raw_data_stat,"enterpriseValue")
        rev=cls.search(raw_data_fin,"annualEbitda")
        evbysales=[]
        years=[]
        for dataPoint in rev:
            date=str(datetime.strptime(dataPoint["asOfDate"],"%Y-%m-%d").year)
            years.append(date)
            evbysales.append(int(ev["raw"]/dataPoint["reportedValue"]["raw"]))
        fig=plt.figure(figsize=(6,4))
        bar=plt.bar(years,evbysales,0.3,color='#132E57')
        yticks=[int(i) for i in evbysales]
        #plt.gca().set_yticks(yticks)
        plt.gca().spines["left"].set_visible(False)
        plt.gca().spines["right"].set_visible(False)
        plt.gca().spines["top"].set_visible(False)
        plt.gca().spines["bottom"].set_color("#DDDDDD")
        for rect in bar:
            height=rect.get_height()
            plt.text(rect.get_x()+0.15,height/2,str(height)+"x",ha='center', va='bottom')
        #plt.gcf().subplots_adjust(top=1,bottom=0.1,left=0.1)
        plt.gcf().subplots_adjust(top=1,bottom=0.25)
        filename=uuid.uuid4().hex[:8].upper()
        path=os.path.join(os.getcwd(),'Graphs')
        if  not os.path.exists(path):
            os.mkdir(path)
        adress=path+'\\'+filename+'.png'
        plt.savefig(adress)
        return adress      
    @classmethod
    def PbyE(cls,raw_data_stat,raw_data_fin):
        ev=cls.search(raw_data_stat,"regularMarketPrice")
        rev=cls.search(raw_data_fin,"annualBasicEPS")
        evbysales=[]
        years=[]
        for dataPoint in rev:
            date=str(datetime.strptime(dataPoint["asOfDate"],"%Y-%m-%d").year)
            years.append(date)
            evbysales.append(int(ev["raw"]/dataPoint["reportedValue"]["raw"]))
        fig=plt.figure(figsize=(6,4))
        bar=plt.bar(years,evbysales,0.3,color='#132E57')
        yticks=[int(i) for i in evbysales]
        #plt.gca().set_yticks(yticks)
        plt.gca().spines["left"].set_visible(False)
        plt.gca().spines["right"].set_visible(False)
        plt.gca().spines["top"].set_visible(False)
        plt.gca().spines["bottom"].set_color("#DDDDDD")
        for rect in bar:
            height=rect.get_height()
            plt.text(rect.get_x()+0.15,height/2,str(height)+"x",ha='center', va='bottom')
        #plt.gcf().subplots_adjust(top=1,bottom=0.1,left=0.1)
        plt.gcf().subplots_adjust(top=1,bottom=0.25)
        filename=uuid.uuid4().hex[:8].upper()
        path=os.path.join(os.getcwd(),'Graphs')
        if  not os.path.exists(path):
            os.mkdir(path)
        adress=path+'\\'+filename+'.png'
        plt.savefig(adress)
        return adress
    @classmethod
    def ppt_output(cls,pr,company,style):
        raw_data_statistics=cls.base_fetch({"symbol":company,"region":"US"},cls.__url__2)
        raw_data_financials=cls.base_fetch({"symbol":company,"region":"US"},cls.__url__1)
        try:
            HEADING_FONT=style["heading-font"]
        except:
            HEADING_FONT="Calibri"
        try:
            HEADING_FONT_SIZE=int(style["heading-font-size"])
        except:
            HEADING_FONT_SIZE=35
        try:
            SUB_HEADING_FONT=style["sub-heading-font"]
        except:
            SUB_HEADING_FONT="Calibri"
        try:
            BODY_TEXT_FONT=style["body-text-font"]
        except:
            BODY_TEXT_FONT="Calibri"
        try:
            BODY_TEXT_FONT_SIZE=style["body-text-font-size"]
        except:
            BODY_TEXT_FONT_SIZE=13
        BODY_TEXT_HEADING_FONT_SIZE=BODY_TEXT_FONT_SIZE+1
        try:
            THEME_COLOR=tuple(int(255*i) for i in Colour.to_rgb(style["theme-color"]))
        except:
            THEME_COLOR=(1,39,99)
        #  Adding slide
        layout=pr.slide_layouts[6] # Blank layout
        slide=pr.slides.add_slide(layout)
        # # Heading
        heading=Heading(HEADING_FONT_SIZE,THEME_COLOR,HEADING_FONT,(Inches(0.2),Inches(0),Inches(8),Inches(0.5)),"Valuation Outputs",slide)
        heading.create_heading()
        # Separation line
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0),Inches(0.7),Inches(2.1),Inches(0.7))
        fill=line.line
        fill.color.rgb=RGBColor(184,25,4)
        # EV/Sales
        heading_shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1),Inches(0.8),Inches(6.5),Inches(0.3))
        heading_txtFrame=heading_shape.text_frame
        heading_txtFrame.text="EV/Sales"
        try:
            heading_txtFrame.fit_text(font_family=str(SUB_HEADING_FONT))
        except:
            heading_txtFrame.fit_text(font_family=u"Calibri")
        shape_fill=heading_shape.fill
        shape_fill.solid()
        shape_fill.fore_color.rgb=RGBColor(*THEME_COLOR)
        # Graph
        img=cls.EVbySales(raw_data_statistics,raw_data_financials)
        graph=slide.shapes.add_picture(img,Inches(1),Inches(1.3),Inches(6.5),Inches(3))
        os.remove(img)
        # EV/EBITDA
        heading_shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(9.5),Inches(0.8),Inches(6.5),Inches(0.3))
        heading_txtFrame=heading_shape.text_frame
        heading_txtFrame.text="EV/EBITDA"
        try:
            heading_txtFrame.fit_text(font_family=str(SUB_HEADING_FONT))
        except:
            heading_txtFrame.fit_text(font_family=u"Calibri")
        shape_fill=heading_shape.fill
        shape_fill.solid()
        shape_fill.fore_color.rgb=RGBColor(*THEME_COLOR)
        # Graph
        img=cls.EVbyEbitda(raw_data_statistics,raw_data_financials)
        graph=slide.shapes.add_picture(img,Inches(9.5),Inches(1.3),Inches(6.5),Inches(3))
        os.remove(img)
        # EV/EBIT
        heading_shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1),Inches(4.3),Inches(6.5),Inches(0.3))
        heading_txtFrame=heading_shape.text_frame
        heading_txtFrame.text="EV/EBIT"
        try:
            heading_txtFrame.fit_text(font_family=str(SUB_HEADING_FONT))
        except:
            heading_txtFrame.fit_text(font_family=u"Calibri")
        shape_fill=heading_shape.fill
        shape_fill.solid()
        shape_fill.fore_color.rgb=RGBColor(*THEME_COLOR)
        # Graph
        img=cls.EVbyEbit(raw_data_statistics,raw_data_financials)
        graph=slide.shapes.add_picture(img,Inches(1),Inches(4.8),Inches(6.5),Inches(3))
        os.remove(img)
        # P/E
        heading_shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(9.5),Inches(4.3),Inches(6.5),Inches(0.3))
        heading_txtFrame=heading_shape.text_frame
        heading_txtFrame.text="P/E"
        try:
            heading_txtFrame.fit_text(font_family=str(SUB_HEADING_FONT))
        except:
            heading_txtFrame.fit_text(font_family=u"Calibri")
        shape_fill=heading_shape.fill
        shape_fill.solid()
        shape_fill.fore_color.rgb=RGBColor(*THEME_COLOR)
        # Graph
        img=cls.PbyE(raw_data_statistics,raw_data_financials)
        graph=slide.shapes.add_picture(img,Inches(9.5),Inches(4.8),Inches(6.5),Inches(3))
        os.remove(img)
        return   