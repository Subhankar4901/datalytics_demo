import json
from Api import BaseAPI
import matplotlib.colors as Colour
from pptx.util import Inches
from utils.create_heading import Heading
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from utils.create_table import Table
class CashflowStatement(BaseAPI):
    _property="CashflowStatement"
    def __init__(self) -> None:
        super().__init__()
    @classmethod
    def output(cls,query:dict):
        return cls.base_output(query,property=cls._property)
    @classmethod
    def _structured_table_data(cls,data:dict,required_items:list):
        structured_data=dict()
        indexes=[]
        for item in required_items:
            for key,val in data.items():
                if item in val:
                    indexes.append(item)
                break
        columns={}
        for key,val in data.items():
            column=[]
            for item in indexes:
                column.append(val[item])
            columns[key]=column
        with open('nameConfig.json','r') as f:
            nameConfig=json.load(f)
        new_indexes=[]
        for item in indexes:
            new_indexes.append(nameConfig[item])
        structured_data["indexes"]=new_indexes
        structured_data["columns"]=columns
        return structured_data
    @classmethod
    def ppt_output(cls,query:dict,style:dict,pr):
        data=cls.output(query)
        required_items=["netIncome","depreciation","changeToNetincome","changeToInventory","changeToLiabilities","changeToAccountReceivables",
                        "changeToOperatingActivities","totalCashFromOperatingActivities","investments","capitalExpenditures","otherCashflowsFromInvestingActivities"
                        ,"totalCashflowsFromInvestingActivities","dividendsPaid","repurchaseOfStock","netBorrowings","issuanceOfStock","otherCashflowsFromFinancingActivities",
                        "totalCashFromFinancingActivities","effectOfExchangeRate","changeInCash","Cashatbeginning","netCash"]
         # Constants for ppt
        SLIDE_HEIGHT=6860000/914400 #Inches
        SLIDE_WIDTH=12190000/914400 #Inches
        try:
            HEADING_FONT=style["heading-font"]
        except:
            HEADING_FONT="Calibri"
        try:
            HEADING_FONT_SIZE=int(style["heading-font-size"])
        except:
            HEADING_FONT_SIZE=35
        try:
            SUB_HEADING_FONT=style["sub-heading-font"]
        except:
            SUB_HEADING_FONT="Calibri"
        try:
            BODY_TEXT_FONT=style["body-text-font"]
        except:
            BODY_TEXT_FONT="Calibri"
        try:
            BODY_TEXT_FONT_SIZE=style["body-text-font-size"]
        except:
            BODY_TEXT_FONT_SIZE=13
        BODY_TEXT_HEADING_FONT_SIZE=BODY_TEXT_FONT_SIZE+1
        try:
            THEME_COLOR=tuple(int(255*i) for i in Colour.to_rgb(style["theme-color"]))
        except:
            THEME_COLOR=(1,39,99)
        #  Adding slide
        layout=pr.slide_layouts[6] # Blank layout
        slide=pr.slides.add_slide(layout)
         # # Heading
        heading=Heading(HEADING_FONT_SIZE,THEME_COLOR,HEADING_FONT,(Inches(0.2),Inches(0),Inches(SLIDE_WIDTH/2),Inches(0.5)),"Cash Flow Statement",slide)
        heading.create_heading()
        # Separation line
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0),Inches(0.7),Inches(2.5),Inches(0.7))
        fill=line.line
        fill.color.rgb=RGBColor(184,25,4)
        column_width={"normal":2,"index":3}
        # Table
        structured_data=cls._structured_table_data(data,required_items)
        table=Table(structured_data,(4,1,12,len(structured_data["indexes"])*0.2),column_width,slide)
        table.create()
        