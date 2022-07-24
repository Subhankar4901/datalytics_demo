import json
import sys
from Api import BaseAPI
import statistics
from numerize import numerize
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE,PP_ALIGN,MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR,MSO_SHAPE
import matplotlib.colors as Colour
from utils.create_table import Table
from  utils.create_heading import Heading
from datetime import datetime
class ValuationSummary(BaseAPI):
    __url__1="https://yh-finance.p.rapidapi.com/stock/v2/get-statistics"
    __url__2="https://yh-finance.p.rapidapi.com/stock/v2/get-financials"
    @classmethod
    def output(cls,companies:list):
        required_items=["regularMarketPrice","marketCap","enterpriseValue","annualTotalRevenue","annualEbitda","annualOperatingIncome","annualBasicEPS","trailingPE"]
        normal=["regularMarketPrice","marketCap","enterpriseValue"]
        quotion=["annualTotalRevenue","annualEbitda","annualOperatingIncome","annualBasicEPS"]
        company_data={}
        for company in companies:
            raw_data_normal=cls.base_fetch({'symbol':company},url=cls.__url__1)
            raw_data_quotion=cls.base_fetch({'symbol':company},url=cls.__url__2)
            nameConfig={
                    "regularMarketPrice":"Stock Price",
                    "marketCap":"Equity Value",
                    "enterpriseValue":"Enterprise Value",
                    "annualTotalRevenue":"EV/Revenue",
                    "annualEbitda":"EV/EBITDA",
                    "annualOperatingIncome":"EV/EBIT",
                    "annualBasicEPS":"Price/Earnings",
                    "trailingPE":"Trailing PE"
            }
            company_data[company]={}
            for item in required_items:
                if item in normal:
                    data=cls.search(raw_data_normal,item)
                    company_data[company][nameConfig[item]]=data["raw"]
                elif item in quotion:
                    data=cls.search(raw_data_quotion,item)
                    for itemData in data:
                        try:
                            year=datetime.strptime(itemData["asOfDate"],"%Y-%m-%d").year
                            key=nameConfig[item]+"\n"+str(year)
                            if itemData["reportedValue"]["raw"]<0:
                                company_data[company][key]="nm"
                            elif item=="annualBasicEPS":
                                company_data[company][key]=company_data[company][nameConfig["regularMarketPrice"]]/itemData["reportedValue"]["raw"]
                            else:
                                company_data[company][key]=company_data[company][nameConfig["marketCap"]]/itemData["reportedValue"]["raw"]
                        except:
                            year+=1
                            key=nameConfig[item]+"\n"+str(year)
                            company_data[company][key]=None
                else:
                    data=cls.search(raw_data_quotion,item)
                    company_data[company][nameConfig[item]]=data["raw"]           
        stats=['High','Median','Mean','Low']
        indexes=companies+stats
        columns={}
        column_names=[i for i in company_data[companies[0]].keys()]
        for column in column_names:
            column_data=[]
            for company in companies:
                try:
                    if column.split("-")[0] in quotion:
                        if company_data[company][column]=="nm" or company_data[company][column]==None: 
                            column_data.append(company_data[company][column])
                        else:    
                            column_data.append(numerize.numerize(company_data[company][column])+"x")
                    else:
                        column_data.append(numerize.numerize(company_data[company][column]))
                except:
                    column_data.append(None)                 
            columns[column]=column_data
        raw_columns={}
        for column in column_names:
            column_data=[]
            for company in companies:
                try:
                    if isinstance(company_data[company][column],str):
                        column_data.append(None)
                    else:
                        column_data.append(company_data[company][column])
                except:
                    column_data.append(None)
            raw_columns[column]=column_data
        for column in raw_columns:
            try:
                High=numerize.numerize(max([val for val in raw_columns[column] if val is not None]),2)
                if column.split("-")[0] in quotion:
                    High+="x"
            except:
                High=None
            try:
                Median=numerize.numerize(statistics.median([val for val in raw_columns[column] if val is not None]),2)
                if column.split("-")[0] in quotion:
                    Median+="x"
            except:
                Median=None
            try:
                Mean=numerize.numerize(statistics.mean([val for val in raw_columns[column] if val is not None]),2)
                if column.split("-")[0] in quotion:
                    Mean+="x"
            except:
                Mean=None
            try:
                Low=numerize.numerize(min([val for val in raw_columns[column] if val is not None]),2)
                if column.split("-")[0] in quotion:
                    Low+="x"
            except:
                Low=None
            columns[column]=columns[column]+[High,Median,Mean,Low]
        return {"indexes":indexes,"columns":columns}
    @classmethod
    def ppt_output(cls,pr,companies,style:dict):
        data=cls.output(companies)
        try:
            HEADING_FONT=style["heading-font"]
        except:
            HEADING_FONT="Calibri"
        try:
            HEADING_FONT_SIZE=int(style["heading-font-size"])
        except:
            HEADING_FONT_SIZE=35
        try:
            SUB_HEADING_FONT=style["sub-heading-font"]
        except:
            SUB_HEADING_FONT="Calibri"
        try:
            BODY_TEXT_FONT=style["body-text-font"]
        except:
            BODY_TEXT_FONT="Calibri"
        try:
            BODY_TEXT_FONT_SIZE=style["body-text-font-size"]
        except:
            BODY_TEXT_FONT_SIZE=13
        BODY_TEXT_HEADING_FONT_SIZE=BODY_TEXT_FONT_SIZE+1
        try:
            THEME_COLOR=tuple(int(255*i) for i in Colour.to_rgb(style["theme-color"]))
        except:
            THEME_COLOR=(1,39,99)
        #  Adding slide
        layout=pr.slide_layouts[6] # Blank layout
        slide=pr.slides.add_slide(layout)
        # # Heading
        heading=Heading(HEADING_FONT_SIZE,THEME_COLOR,HEADING_FONT,(Inches(0.2),Inches(0),Inches(8),Inches(0.5)),"Valuation Summary",slide)
        heading.create_heading()
        # Separation line
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0),Inches(0.7),Inches(2.1),Inches(0.7))
        fill=line.line
        fill.color.rgb=RGBColor(184,25,4)
        # Table
        column_width={"normal":0.8
                      ,"index":0.8}
        table=Table(data,(0.6,1.7,len(data["columns"])*0.9+0.8,len(data["indexes"])*0.5),column_width,slide)
        table.create()
        return