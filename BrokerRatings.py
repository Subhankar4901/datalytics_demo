import matplotlib.pyplot as plt
import numpy as np
from sympy import false
from Api import BaseAPI
import statistics
from numerize import numerize
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE,PP_ALIGN,MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR,MSO_SHAPE
import matplotlib.colors as Colour
from utils.create_table import Table
from  utils.create_heading import Heading
from datetime import datetime
import os
import uuid
class BrokerRatings(BaseAPI):
    __url__="https://yh-finance.p.rapidapi.com/stock/v2/get-analysis"
    @classmethod
    def formatGrade(cls,grade):
        if grade in ["Buy","Overweight","Outperform", "Strong Buy","Market Outperform"]:
            return "Buy"
        if  grade in ["Neutral","Hold","Peer Perform","Equal-Weight","Market Perform", "In-Line","Mixed","Sector Perform"]:
            return "Hold"
        if grade in ["Sell","Underperform","Under weight", "Reduce","Underweight"]:
            return "Sell"
        return "Unrecognised"
    @classmethod
    def table_output(cls,raw_data):
        data=cls.search(raw_data,"upgradeDowngradeHistory")
        count=0
        column_data={}
        column_data["Broker"]=[]
        column_data["Rating"]=[]
        indexes=[]
        for dataPoint in data["history"][:15]:
            date=datetime.fromtimestamp(dataPoint["epochGradeDate"])
            datestr=date.strftime("%d-%b-%y")
            grade=cls.formatGrade(dataPoint["toGrade"])
            firm=dataPoint["firm"]
            indexes.append(datestr)
            column_data["Broker"].append(firm)
            column_data["Rating"].append(grade)
        return {"columns":column_data,"indexes":indexes}
            
                  
    @classmethod
    def get_graph(cls,raw_data):
        data=cls.search(raw_data,"upgradeDowngradeHistory")
        count=0
        graph_data={}
        un=[]
        for dataPoint in data["history"]:
            date=datetime.fromtimestamp(dataPoint["epochGradeDate"])
            datestr=date.strftime("%b,%y")
            grade=cls.formatGrade(dataPoint["toGrade"])
            if grade=="Unrecognised":
                un.append(dataPoint["toGrade"])
            if datestr not in graph_data:
                graph_data[datestr]={}
                count+=1
            if grade not in graph_data[datestr]:
                graph_data[datestr][grade]=0
            graph_data[datestr][grade]+=1
            if count==15:
                break
        dates=[]
        buy=[]
        sell=[]
        hold=[]
        unrecognised=[]
        for key,val in graph_data.items():
            dates.append(key)
            if "Buy" in val:
                buy.append(val["Buy"])
            else:
                buy.append(0)
            if "Hold" in val:
                hold.append(val["Hold"])
            else:
                hold.append(0)
            if "Sell" in val:
                sell.append(val["Sell"])
            else:
                sell.append(0)
            if "Unrecognised" in val:
                unrecognised.append(val["Unrecognised"])
            else:
                unrecognised.append(0)
        print(un)
        dates=np.array(dates)
        buy=np.array(buy)
        sell=np.array(sell)
        hold=np.array(hold)
        fig=plt.figure(figsize=(10,6))
        bar1=plt.bar(dates,sell,width=0.3,color="#ff333a")
        bar2=plt.bar(dates,hold,0.3,bottom=sell,color="#ffdd48")
        bar3=plt.bar(dates,buy,0.3,bottom=sell+hold,color="#00c073")
        plt.gca().spines["left"].set_visible(False)
        plt.gca().spines["right"].set_visible(False)
        plt.gca().spines["top"].set_visible(False)
        plt.gca().spines["bottom"].set_color("#DDDDDD")
        plt.gca().tick_params(left=False)
        yticks=set()
        for tick in sell:
            yticks.add(tick)
        for tick in sell+hold:
            yticks.add(tick)
        for tick in sell+hold+buy:
            yticks.add(tick)
        plt.gca().set_yticks(list(yticks))
        plt.gcf().autofmt_xdate()
        plt.gca().set_axisbelow(True)
        plt.gca().yaxis.grid(True,color="#EEEEEE")
        plt.gca().xaxis.grid(False)
        for tick in plt.gca().xaxis.get_major_ticks():
            tick.label.set_fontsize(9)
        for tick in plt.gca().yaxis.get_major_ticks():
            tick.label.set_fontsize(9)
        plt.legend(["Sell","Hold","Buy"])
        plt.gcf().subplots_adjust(top=1,bottom=0.1,left=0.1)
        filename=uuid.uuid4().hex[:8].upper()
        path=os.path.join(os.getcwd(),'Graphs')
        if  not os.path.exists(path):
            os.mkdir(path)
        adress=path+'\\'+filename+'.png'
        plt.savefig(adress)
        return adress
        
    @classmethod
    def ppt_output(cls,pr,company,style):
        try:
            HEADING_FONT=style["heading-font"]
        except:
            HEADING_FONT="Calibri"
        try:
            HEADING_FONT_SIZE=int(style["heading-font-size"])
        except:
            HEADING_FONT_SIZE=35
        try:
            SUB_HEADING_FONT=style["sub-heading-font"]
        except:
            SUB_HEADING_FONT="Calibri"
        try:
            BODY_TEXT_FONT=style["body-text-font"]
        except:
            BODY_TEXT_FONT="Calibri"
        try:
            BODY_TEXT_FONT_SIZE=style["body-text-font-size"]
        except:
            BODY_TEXT_FONT_SIZE=13
        BODY_TEXT_HEADING_FONT_SIZE=BODY_TEXT_FONT_SIZE+1
        try:
            THEME_COLOR=tuple(int(255*i) for i in Colour.to_rgb(style["theme-color"]))
        except:
            THEME_COLOR=(1,39,99)
        #  Adding slide
        layout=pr.slide_layouts[6] # Blank layout
        slide=pr.slides.add_slide(layout)
        # # Heading
        heading=Heading(HEADING_FONT_SIZE,THEME_COLOR,HEADING_FONT,(Inches(0.2),Inches(0),Inches(8),Inches(0.5)),"Broker Ratings",slide)
        heading.create_heading()
        # Separation line
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0),Inches(0.7),Inches(2.1),Inches(0.7))
        fill=line.line
        fill.color.rgb=RGBColor(184,25,4)
        data=cls.base_fetch({"symbol":company,"region":"US"},cls.__url__)
        graph_data=cls.get_graph(data)
        table_data=cls.table_output(data)
        column_width={"normal":2,"index":2}
        table=Table(table_data,(1,1.6,5,len(table_data["indexes"])*0.4),column_width,slide)
        table.create()
        graph=slide.shapes.add_picture(graph_data,Inches(8),Inches(1.5),Inches(10),Inches(6))
        os.remove(graph_data)
        
        
            
            
        
            
                
                
                