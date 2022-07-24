from Api import BaseAPI
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE,PP_ALIGN,MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR,MSO_SHAPE
import matplotlib.colors as Colour
import json
from utils.create_table import Table
from  utils.create_heading import Heading
class BalanceSheet(BaseAPI):
    _property="BalanceSheet"
    __url__="https://yh-finance.p.rapidapi.com/stock/v2/get-balance-sheet"
    def __init__(self) -> None:
        super().__init__()
    
    @classmethod
    def output(cls,query:dict):
        return cls.base_output(query,property=cls._property)    
    @classmethod
    def _structured_table_data(cls,data:dict,required_items:list):
        structured_data=dict()
        indexes=[]
        for item in required_items:
            for key,val in data.items():
                if item in val:
                    indexes.append(item)
                break
        columns={}
        for key,val in data.items():
            column=[]
            for item in indexes:
                column.append(val[item])
            columns[key]=column
        with open('nameConfig.json','r') as f:
            nameConfig=json.load(f)
        new_indexes=[]
        for item in indexes:
            new_indexes.append(nameConfig[item])
        structured_data["indexes"]=new_indexes
        structured_data["columns"]=columns
        return structured_data        
    @classmethod
    def ppt_output(cls,query:dict,style:dict,pr):
        data=cls.output(query)
        items_for_total_current_assets=["cash","shortTermInvestments","netReceivables","inventory","otherCurrentAssets","totalCurrentAssets"]
        items_for_total_assets=["totalCurrentAssets","intangibleAssets","propertyPlantEquipment","netTangibleAssets","goodWill","longTermInvestments",
                        "otherAssets","totalAssets"]
        items_for_total_current_liabilities=["accountsPayable","shortLongTermDebt","otherCurrentLiab","totalCurrentLiabilities"]
        items_for_total_liabilities=["totalCurrentLiabilities","longTermDebt","deferredLongTermAssetCharges","deferredLongTermLiab","minorityInterest","otherLiab","totalLiab"]
        items_for_total_stockholders_equity=["commonStock","otherStockholderEquity","treasuryStock","retainedEarnings","capitalSurplus","totalStockholderEquity"]
        # Constants for ppt
        SLIDE_HEIGHT=6860000/914400 #Inches
        SLIDE_WIDTH=12190000/914400 #Inches
        try:
            HEADING_FONT=style["heading-font"]
        except:
            HEADING_FONT="Calibri"
        try:
            HEADING_FONT_SIZE=int(style["heading-font-size"])
        except:
            HEADING_FONT_SIZE=35
        try:
            SUB_HEADING_FONT=style["sub-heading-font"]
        except:
            SUB_HEADING_FONT="Calibri"
        try:
            BODY_TEXT_FONT=style["body-text-font"]
        except:
            BODY_TEXT_FONT="Calibri"
        try:
            BODY_TEXT_FONT_SIZE=style["body-text-font-size"]
        except:
            BODY_TEXT_FONT_SIZE=13
        BODY_TEXT_HEADING_FONT_SIZE=BODY_TEXT_FONT_SIZE+1
        try:
            THEME_COLOR=tuple(int(255*i) for i in Colour.to_rgb(style["theme-color"]))
        except:
            THEME_COLOR=(1,39,99)
        #  Adding slide
        layout=pr.slide_layouts[6] # Blank layout
        slide=pr.slides.add_slide(layout)
        
        # # Heading
        heading=Heading(HEADING_FONT_SIZE,THEME_COLOR,HEADING_FONT,(Inches(0.2),Inches(0),Inches(SLIDE_WIDTH/2),Inches(0.5)),"Balance Sheet",slide)
        heading.create_heading()
        # Separation line
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0),Inches(0.7),Inches(2.1),Inches(0.7))
        fill=line.line
        fill.color.rgb=RGBColor(184,25,4)
        
        # Sub heading
        assets_heading=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(3.5),Inches(1.1),Inches(0.9),Inches(0.4))
        assets_heading.text_frame.text="Assets"
        assets_heading.fill.solid()
        assets_heading.fill.fore_color.rgb=RGBColor(17, 237, 35)
        
        liab_heading=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(13),Inches(1.1),Inches(1.3),Inches(0.4))
        liab_heading.text_frame.text="Liabilities"
        liab_heading.fill.solid()
        liab_heading.fill.fore_color.rgb=RGBColor(245, 66, 66)
        
        column_width={"normal":1,"index":3}
        # Total current assets table.
        total_current_assets_data=cls._structured_table_data(data,items_for_total_current_assets)
        table=Table(total_current_assets_data,(1,1.7,6,len(total_current_assets_data["indexes"])*0.2),column_width,slide)
        table.create()
        #  Total assets table.
        total_assets_data=cls._structured_table_data(data,items_for_total_assets)
        table=Table(total_assets_data,(1,4,6,len(total_assets_data["indexes"])*0.2),column_width,slide)
        table.create()
        # Total current liabilities table.
        total_current_liabilities_data=cls._structured_table_data(data,items_for_total_current_liabilities)
        table=Table(total_current_liabilities_data,(10,1.7,6,len(total_current_liabilities_data["indexes"])*0.2),column_width,slide)
        table.create()
        # Total liabilities table.
        total_liabilities_data=cls._structured_table_data(data,items_for_total_liabilities)
        table=Table(total_liabilities_data,(10,3.5,6,len(total_liabilities_data["indexes"])*0.2),column_width,slide)
        table.create()
        # Total stockholder's equity table.
        total_stockholders_equity_data=cls._structured_table_data(data,items_for_total_stockholders_equity)
        table=Table(total_stockholders_equity_data,(10,5.3,6,len(total_stockholders_equity_data["indexes"])*0.2),column_width,slide)
        table.create()
        return       
        
                
        
            
        
        

        
        
        
        
            
        