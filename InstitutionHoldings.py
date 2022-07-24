from Api import BaseAPI
import statistics
from numerize import numerize
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE,PP_ALIGN,MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR,MSO_SHAPE
import matplotlib.colors as Colour
from utils.create_table import Table
from  utils.create_heading import Heading
from datetime import datetime
class InstitutionalHoldings(BaseAPI):
    __url__1="https://yh-finance.p.rapidapi.com/stock/v2/get-holders"
    @classmethod
    def output(cls,raw_data):
        data=cls.search(raw_data,"institutionOwnership")
        column_data={}
        indexes=[]
        column_data["Reported Date"]=[]
        column_data["% Owned"]=[]
        column_data["No. of shares"]=[]
        for dataPoint in data["ownershipList"]:
            indexes.append(dataPoint["organization"])
            date=datetime.strptime(dataPoint["reportDate"]["fmt"],"%Y-%m-%d")
            column_data["Reported Date"].append(date.strftime("%d %b,%Y"))
            column_data["% Owned"].append(dataPoint["pctHeld"]["fmt"])
            column_data["No. of shares"].append(dataPoint["position"]["longFmt"])
        return {"indexes":indexes,"columns":column_data}
    @classmethod
    def ppt_output(cls,pr,company,style):
        try:
            HEADING_FONT=style["heading-font"]
        except:
            HEADING_FONT="Calibri"
        try:
            HEADING_FONT_SIZE=int(style["heading-font-size"])
        except:
            HEADING_FONT_SIZE=35
        try:
            SUB_HEADING_FONT=style["sub-heading-font"]
        except:
            SUB_HEADING_FONT="Calibri"
        try:
            BODY_TEXT_FONT=style["body-text-font"]
        except:
            BODY_TEXT_FONT="Calibri"
        try:
            BODY_TEXT_FONT_SIZE=style["body-text-font-size"]
        except:
            BODY_TEXT_FONT_SIZE=13
        BODY_TEXT_HEADING_FONT_SIZE=BODY_TEXT_FONT_SIZE+1
        try:
            THEME_COLOR=tuple(int(255*i) for i in Colour.to_rgb(style["theme-color"]))
        except:
            THEME_COLOR=(1,39,99)
        #  Adding slide
        layout=pr.slide_layouts[6] # Blank layout
        slide=pr.slides.add_slide(layout)
        # # Heading
        heading=Heading(HEADING_FONT_SIZE,THEME_COLOR,HEADING_FONT,(Inches(0.2),Inches(0),Inches(8),Inches(0.5)),"Institutional Holdings",slide)
        heading.create_heading()
        # Separation line
        line=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(0),Inches(0.7),Inches(2.1),Inches(0.7))
        fill=line.line
        fill.color.rgb=RGBColor(184,25,4)
        raw_data=cls.base_fetch({"symbol":company,"region":"US"},cls.__url__1)
        data=cls.output(raw_data)
        column_width={"normal":2,"index":3}
        table=Table(data,(4,1.5,12,len(data["indexes"])*0.4),column_width,slide)
        table.create()
        return